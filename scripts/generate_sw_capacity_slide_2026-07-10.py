#!/usr/bin/env python3
"""Single-slide PPT: Sandwich UK capacity building summary for R&D review (Jul 2026)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path("/workspace/汇报/SW研发述职_产能建设一页_2026-07-10.pptx")

NAVY = RGBColor(0x0F, 0x2B, 0x46)
TEAL = RGBColor(0x00, 0x96, 0x88)
ACCENT = RGBColor(0xC9, 0xA2, 0x27)
GREY = RGBColor(0x5A, 0x6A, 0x7A)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x2C, 0x3E, 0x50)

FONT = "Microsoft YaHei"


def _fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color


def _textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _set_para(p, text, *, size=11, bold=False, color=TEXT, align=PP_ALIGN.LEFT, space_after=4):
    p.text = text
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(space_after)


def _add_bullets(tf, items, *, base_size=10.5):
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, (label, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.font.name = FONT
        p.font.size = Pt(base_size)
        p.space_after = Pt(7)
        if label:
            r0 = p.add_run()
            r0.text = f"{label}："
            r0.font.bold = True
            r0.font.color.rgb = NAVY
            r1 = p.add_run()
            r1.text = body
            r1.font.bold = False
            r1.font.color.rgb = TEXT
        else:
            r = p.add_run()
            r.text = body
            r.font.color.rgb = TEXT


def _card_header(slide, left, top, width, title, accent: RGBColor):
    hdr_h = Inches(0.38)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.06), hdr_h)
    _fill(bar, accent)
    bar.line.fill.background()
    box = _textbox(slide, left + Inches(0.12), top, width - Inches(0.12), hdr_h)
    _set_para(box.text_frame.paragraphs[0], title, size=13, bold=True, color=NAVY)


def _draw_layout_schematic(slide, left, top, width, height):
    """Simple block layout placeholder when no render image is available."""
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(outer, WHITE)
    outer.line.color.rgb = RGBColor(0xD4, 0xDD, 0xE6)

    title = _textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.3))
    _set_para(title.text_frame.paragraphs[0], "园区布局示意（概念阶段）", size=11, bold=True, color=NAVY)

    # Existing B902 block
    exist = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), top + Inches(0.55), Inches(1.55), Inches(1.05))
    _fill(exist, RGBColor(0xE8, 0xEE, 0xF4))
    exist.line.color.rgb = RGBColor(0x9A, 0xA8, 0xB6)
    tb = _textbox(slide, left + Inches(0.25), top + Inches(0.72), Inches(1.25), Inches(0.7))
    tf = tb.text_frame
    _set_para(tf.paragraphs[0], "既有 B902", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    p2 = tf.add_paragraph()
    _set_para(p2, "改造区\n（色谱/冻干/密闭升级）", size=8, color=GREY, align=PP_ALIGN.CENTER)

    # Extension block
    ext = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(1.85), top + Inches(0.55), Inches(1.15), Inches(1.05))
    _fill(ext, RGBColor(0xD9, 0xEA, 0xE8))
    ext.line.color.rgb = TEAL
    tb2 = _textbox(slide, left + Inches(1.7), top + Inches(0.78), Inches(0.95), Inches(0.55))
    tf2 = tb2.text_frame
    _set_para(tf2.paragraphs[0], "东侧扩建", size=9, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    p3 = tf2.add_paragraph()
    _set_para(p3, "反应/加氢/过滤", size=8, color=GREY, align=PP_ALIGN.CENTER)

    # HP lab block
    hp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.2), top + Inches(1.75), Inches(2.8), Inches(0.55))
    _fill(hp, RGBColor(0xFB, 0xF3, 0xE0))
    hp.line.color.rgb = ACCENT
    tb3 = _textbox(slide, left + Inches(0.3), top + Inches(1.88), Inches(2.3), Inches(0.35))
    _set_para(tb3.text_frame.paragraphs[0], "高活实验室（已选空间 · 布局初定）", size=9, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    note = _textbox(slide, left + Inches(0.15), top + height - Inches(0.38), width - Inches(0.3), Inches(0.3))
    _set_para(
        note.text_frame.paragraphs[0],
        "可替换为实际效果图或正式平面布置图",
        size=8,
        color=GREY,
        align=PP_ALIGN.CENTER,
    )


def build_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1))
    _fill(bar, NAVY)
    bar.line.fill.background()

    # Title
    title = _textbox(slide, Inches(0.45), Inches(0.28), Inches(8.5), Inches(0.55))
    _set_para(title.text_frame.paragraphs[0], "英国三明治站点 · 产能建设进展", size=26, bold=True, color=NAVY)
    sub = _textbox(slide, Inches(0.45), Inches(0.82), Inches(8), Inches(0.35))
    _set_para(sub.text_frame.paragraphs[0], "研发支撑能力升级  |  2026年7月", size=12, color=GREY)

    # Three content cards
    card_top = Inches(1.28)
    card_h = Inches(4.55)
    gap = Inches(0.12)
    card_w = Inches(3.35)
    x0 = Inches(0.42)

    cards = [
        (
            "B902 东侧扩建",
            "新增规模化产能",
            TEAL,
            [
                ("定位", "反应、加氢、过滤干燥能力提升，支撑英国中长期产能"),
                ("范围", "东侧新建四层主体及设备夹层，总建筑面积约 3,099 平方米；10 台反应釜、2,500 升加氢釜、3 套过滤干燥机"),
                ("投资", "可行性阶段投资测算已完成"),
                ("进展", "可行性研究已完成（2026年5月）；后续进入概念设计及详细设计"),
                ("节点", "规划许可 → 施工建设 → 2030年5月竣工目标"),
            ],
        ),
        (
            "B902 既有厂房改造",
            "近期投产能力",
            RGBColor(0x5B, 0x6E, 0xAE),
            [
                ("定位", "在现有厂房内补齐制备液相色谱、冻干及五级密闭日常运行能力"),
                ("制备液相色谱 + 冻干", "技术可行性研究已完成；投资测算已完成；冻干长周期驱动分期实施"),
                ("C1 模块五级密闭升级", "范围与费用测算已完成；须与液相色谱改造同步交付"),
                ("", "改造线优先兑现 2027 年可运行能力；与扩建线并行、分阶段投产"),
            ],
        ),
        (
            "高活实验室",
            "新增高活运行能力",
            ACCENT,
            [
                ("定位", "新增高活实验与运行能力，与五级密闭升级及制备液相色谱能力衔接"),
                ("投资", "投资测算同步推进中"),
                (
                    "进展",
                    "已与三明治站点、发现园进行多轮接洽；完成布局初步合理设计与设备数量测算，"
                    "初步确认可行性，选定空间并完成设计准备",
                ),
            ],
        ),
    ]

    for i, (title_txt, subtitle_txt, accent, bullets) in enumerate(cards):
        left = x0 + i * (card_w + gap)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, card_top, card_w, card_h)
        _fill(card, WHITE)
        card.line.color.rgb = RGBColor(0xE0, 0xE6, 0xEC)

        _card_header(slide, left + Inches(0.14), card_top + Inches(0.14), card_w - Inches(0.28), title_txt, accent)
        st = _textbox(slide, left + Inches(0.14), card_top + Inches(0.5), card_w - Inches(0.28), Inches(0.28))
        _set_para(st.text_frame.paragraphs[0], subtitle_txt, size=10, color=GREY)

        body = _textbox(slide, left + Inches(0.14), card_top + Inches(0.82), card_w - Inches(0.28), card_h - Inches(0.95))
        _add_bullets(body.text_frame, bullets)

    # Right layout schematic
    layout_left = Inches(10.95)
    layout_top = Inches(1.28)
    layout_w = Inches(2.0)
    layout_h = Inches(4.55)
    _draw_layout_schematic(slide, layout_left, layout_top, layout_w, layout_h)

    # Bottom KPI strip (no timeline bar)
    kpi_top = Inches(6.0)
    kpi_w = Inches(4.05)
    kpis = [
        ("可行性研究", "扩建 / 改造 已完成", TEAL),
        ("投资测算", "扩建 / 改造 已完成", NAVY),
        ("高活实验室", "布局与空间已初步确定", ACCENT),
    ]
    for i, (k, v, c) in enumerate(kpis):
        left = Inches(0.42) + i * (kpi_w + Inches(0.15))
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, kpi_top, kpi_w, Inches(0.72))
        _fill(box, LIGHT_BG)
        box.line.color.rgb = c
        tb = _textbox(slide, left + Inches(0.18), kpi_top + Inches(0.1), kpi_w - Inches(0.36), Inches(0.55))
        tf = tb.text_frame
        _set_para(tf.paragraphs[0], k, size=11, bold=True, color=NAVY)
        _set_para(tf.add_paragraph(), v, size=10, color=GREY)

    return prs


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = build_slide()
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
