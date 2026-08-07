#!/usr/bin/env python3
"""Presentation-grade flowchart: SW B902 purified water sampling workflow."""

from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch, Rectangle

OUT_PNG = Path("/workspace/汇报/纯化水取样/SW_纯化水取样流程图.png")
OUT_PPTX = Path("/workspace/汇报/纯化水取样/SW_纯化水取样流程_一页.pptx")

NAVY = "#0f2b46"
BLUE = "#1a4a6e"
TEAL = "#009688"
ACCENT = "#c9a227"
WARN = "#b43a2a"
GREY = "#5a6a7a"
LIGHT = "#f4f6f8"
CARD = "#ffffff"
EDGE = "#dde3e8"


def _font() -> str:
    for name in ("Microsoft YaHei", "WenQuanYi Micro Hei", "Droid Sans Fallback"):
        if any(f.name == name for f in font_manager.fontManager.ttflist):
            return name
    return "DejaVu Sans"


def _box(ax, x, y, w, h, *, face=CARD, edge=EDGE, lw=1.0, radius=0.012, z=2):
    p = FancyBboxPatch(
        (x, y), w, h,
        boxstyle=f"round,pad=0,rounding_size={radius}",
        facecolor=face, edgecolor=edge, linewidth=lw, zorder=z,
    )
    ax.add_patch(p)
    return p


def _text(ax, x, y, s, *, size=8, color=NAVY, ha="left", va="top", weight="normal", z=5):
    ax.text(x, y, s, fontsize=size, color=color, ha=ha, va=va,
            fontweight=weight, zorder=z, linespacing=1.5)


def _arrow(ax, x1, y1, x2, y2, *, color=BLUE, lw=1.6, style="-|>", ms=9, ls="-", z=3):
    ax.add_patch(FancyArrowPatch(
        (x1, y1), (x2, y2), arrowstyle=f"{style},head_width=3.2,head_length={ms * 0.42}",
        color=color, linewidth=lw, linestyle=ls, zorder=z,
        shrinkA=0, shrinkB=0, mutation_scale=1,
    ))


def _stage(ax, x, y, w, h, idx, title_cn, title_en, bullets, accent):
    """Numbered stage card with header strip and bullet body."""
    _box(ax, x, y, w, h, face=CARD, edge=EDGE, lw=1.1)
    hdr_h = 0.052
    ax.add_patch(Rectangle((x, y + h - hdr_h), w, hdr_h, facecolor=accent,
                           edgecolor="none", zorder=3))
    # number badge
    ax.add_patch(Rectangle((x + 0.006, y + h - hdr_h + 0.008), 0.019, hdr_h - 0.016,
                           facecolor="#ffffff", edgecolor="none", zorder=4))
    _text(ax, x + 0.0155, y + h - hdr_h / 2, str(idx), size=9, color=accent,
          ha="center", va="center", weight="bold")
    _text(ax, x + 0.032, y + h - hdr_h / 2, title_cn, size=9.2, color="#ffffff",
          va="center", weight="bold")
    _text(ax, x + w - 0.008, y + h - hdr_h / 2, title_en, size=6.4,
          color="#e8eef4", ha="right", va="center")

    ty = y + h - hdr_h - 0.022
    for b in bullets:
        _text(ax, x + 0.012, ty, b, size=7.0, color="#2c3e50")
        ty -= 0.038 * (1 + b.count("\n"))


def build_png():
    fam = _font()
    plt.rcParams["font.family"] = fam
    plt.rcParams["axes.unicode_minus"] = False

    fig = plt.figure(figsize=(13.333, 7.5), dpi=150)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    ax.add_patch(Rectangle((0, 0), 1, 1, facecolor=LIGHT, edgecolor="none", zorder=0))

    # ---------- Header ----------
    ax.add_patch(Rectangle((0, 0.975), 1, 0.025, facecolor=NAVY, edgecolor="none", zorder=1))
    _text(ax, 0.028, 0.955, "SW · B902 中试厂房纯化水取样流程", size=17, color=NAVY, weight="bold")
    _text(ax, 0.028, 0.913,
          "Purified Water Sampling Workflow  ·  SOP-BLD-033 + 排程表 QF13070  ·  取样由 B902 运行团队执行，检测全部外委",
          size=8.4, color=GREY)

    # header KPI chips
    chips = [
        ("每周五上午", "Weekly · Friday AM"),
        ("约 8 个点位", "~8 sample points"),
        ("24h 冷链外委", "Sheffield lab"),
        ("48h 内判定", "Review within 48h"),
    ]
    cx = 0.575
    for cn, en in chips:
        _box(ax, cx, 0.905, 0.100, 0.055, face="#e7eefb", edge="#e7eefb")
        _text(ax, cx + 0.050, 0.947, cn, size=8.2, color=NAVY, ha="center", va="center", weight="bold")
        _text(ax, cx + 0.050, 0.921, en, size=5.9, color=GREY, ha="center", va="center")
        cx += 0.106

    # ---------- Main 5-stage flow ----------
    sy, sh = 0.560, 0.310
    sw = 0.174
    gap = 0.0245
    x0 = 0.028
    stages = [
        (1, "取样准备", "Preparation", [
            "• 查 QF13070 确定本周点位与检项",
            "• 标签：楼号/点位/日期/取样人",
            "• 备瓶：Sterilin 500 mL、Nalgene、\n   TOC vial（仅离线时）",
            "• PPE + 一次性丁腈手套",
        ], NAVY),
        (2, "现场取样", "On-site sampling", [
            "• 阀口 70% IPA 喷洒，静置 ≥2 min",
            "• 开阀冲洗 ≥1 min，避免回溅",
            "• 顶空：微生物 ~3 cm / 内毒素 ~1 cm",
            "• 关阀后再次 IPA 喷洒；每区换手套",
        ], BLUE),
        (3, "送检外委", "Dispatch to lab", [
            "• 24 h 内冷链快递至 Sheffield",
            "  合作第三方实验室（UK）",
            "• 周一 出初步结果",
            "• 周四 出正式报告",
        ], TEAL),
        (4, "结果判定", "Result review", [
            "• 收到报告后 48 h 内完成审核",
            "• 核对报告头信息与实验室批准",
            "• 比对内控警戒限 / 行动限",
            "• 评估实验室偏差对水质的影响",
        ], ACCENT),
        (5, "趋势与报告", "Trending & report", [
            "• 每周与上周结果比对，识别重复超限",
            "• 年度趋势报告，确认撬块运行满意",
            "• 经中试代表 + QA 审批",
            "• 归档文件管理系统（DMS）",
        ], "#5b6eae"),
    ]
    xs = []
    for i, (idx, cn, en, bullets, color) in enumerate(stages):
        x = x0 + i * (sw + gap)
        xs.append(x)
        _stage(ax, x, sy, sw, sh, idx, cn, en, bullets, color)
        if i:
            _arrow(ax, x - gap + 0.002, sy + sh / 2, x - 0.003, sy + sh / 2, color=GREY, lw=1.8)

    # repeat-sample feedback loop (stage 4 -> stage 2)
    loop_y = sy - 0.038
    x_from = xs[3] + sw / 2
    x_to = xs[1] + sw / 2
    _arrow(ax, x_from, sy + 0.002, x_from, loop_y, color=WARN, lw=1.2, ls=(0, (4, 2)), ms=7)
    ax.plot([x_to, x_from], [loop_y, loop_y], color=WARN, lw=1.2, ls=(0, (4, 2)), zorder=3)
    _arrow(ax, x_to, loop_y, x_to, sy - 0.002, color=WARN, lw=1.2, ls=(0, (4, 2)), ms=7)
    _text(ax, (x_from + x_to) / 2, loop_y - 0.006, "必要时下一个周五对该点复测  Repeat sample if required",
          size=6.6, color=WARN, ha="center", va="top")

    # ---------- Lower left: sampling order ----------
    ly, lh = 0.290, 0.210
    _box(ax, 0.028, ly, 0.470, lh, face=CARD, edge=EDGE, lw=1.1)
    _text(ax, 0.042, ly + lh - 0.020, "取样顺序：由优到劣  Best → worst water quality",
          size=9, color=NAVY, weight="bold")

    order = [
        ("用水点", "User points", TEAL),
        ("SP12", "分配总管", TEAL),
        ("SP09", "UV 后", BLUE),
        ("CEDI-V103", "发生撬", ACCENT),
        ("SP06", "软化水", GREY),
    ]
    bw, bh = 0.079, 0.052
    bx = 0.042
    by = ly + 0.092
    for i, (cn, sub, color) in enumerate(order):
        _box(ax, bx, by, bw, bh, face="#f8fafb", edge=color, lw=1.3)
        _text(ax, bx + bw / 2, by + bh - 0.014, cn, size=7.8, color=NAVY,
              ha="center", va="center", weight="bold")
        _text(ax, bx + bw / 2, by + 0.015, sub, size=6.2, color=GREY, ha="center", va="center")
        if i < len(order) - 1:
            _arrow(ax, bx + bw + 0.002, by + bh / 2, bx + bw + 0.014, by + bh / 2,
                   color=GREY, lw=1.3, ms=7)
        bx += bw + 0.016

    _text(ax, 0.042, ly + 0.072,
          "轮换机制：SP09 / CEDI-V103 每周必取；SP12 每周微生物并按周期叠加内毒素、化学；SP06 约每三周（仅参考 FIO）；\n"
          "各车间用水点（GP1 / GP2 / FG1 / C 区反应釜、干燥机、转料站）按 18 周排程轮换逐点覆盖。",
          size=6.6, color=GREY)

    # ---------- Lower middle: limits ----------
    _box(ax, 0.512, ly, 0.235, lh, face=CARD, edge=EDGE, lw=1.1)
    _text(ax, 0.526, ly + lh - 0.020, "内控限度  In-house limits",
          size=9, color=NAVY, weight="bold")
    rows = [
        ("点位 Point", "警戒 Alert", "行动 Action"),
        ("SP12 / 用水点", "5", "10"),
        ("SP09", "5", "10"),
        ("CEDI-V103", "30", "50"),
        ("SP06（软化水）", "FIO", "FIO"),
    ]
    ry = ly + lh - 0.064
    for i, (a, b, c) in enumerate(rows):
        hdr = i == 0
        if hdr:
            ax.add_patch(Rectangle((0.522, ry - 0.004), 0.215, 0.024,
                                   facecolor=NAVY, edgecolor="none", zorder=3))
        col = "#ffffff" if hdr else "#2c3e50"
        _text(ax, 0.528, ry + 0.008, a, size=6.4, color=col, va="center",
              weight="bold" if hdr else "normal")
        _text(ax, 0.671, ry + 0.008, b, size=6.0 if hdr else 6.4, color=col,
              ha="center", va="center", weight="bold" if hdr else "normal")
        _text(ax, 0.716, ry + 0.008, c, size=6.0 if hdr else 6.4, color=col,
              ha="center", va="center", weight="bold" if hdr else "normal")
        ry -= 0.024
    _text(ax, 0.526, ly + 0.042,
          "单位 cfu/100 mL；SP12 与用水点铜绿假单胞菌不得检出。\n内毒素限 0.25 EU/mL（无菌 / 注射级项目按需增加）。",
          size=6.2, color=GREY)

    # ---------- Lower right: online monitoring ----------
    _box(ax, 0.760, ly, 0.212, lh, face="#fffdf5", edge=ACCENT, lw=1.3)
    _text(ax, 0.774, ly + lh - 0.020, "在线监测  Online monitoring",
          size=9, color=NAVY, weight="bold")
    _text(ax, 0.774, ly + lh - 0.054,
          "• 分配总管电导率与 TOC 为在线连续监测\n"
          "• 仅当在线仪表异常，或电导率超\n"
          "   USP Stage 1 时，才另行离线取样\n"
          "• 取样前查撬块控制盘：软化器\n"
          "   REGEN 灯亮时暂不取样",
          size=6.6, color="#2c3e50")

    # ---------- Bottom: disposition branches ----------
    dy, dh = 0.072, 0.172
    _text(ax, 0.028, dy + dh + 0.020, "判定与处置  Disposition",
          size=9.5, color=NAVY, weight="bold")
    branches = [
        ("合格 In-spec", TEAL, [
            "• 记录并纳入每周趋势比对",
            "• 按原排程继续下周取样",
        ]),
        ("超警戒限 Alert exceeded", ACCENT, [
            "• 分配系统连续两周超限，或同一周内两个样品超限",
            "• 通知 QA，由用户、设施/运行团队与 QA 会商处置",
        ]),
        ("超行动限 Action exceeded", WARN, [
            "• 立即通知运行与设施负责人及 QA 进行影响评估",
            "• 按 SOP-QAC-002 立偏差；处置记录留存备审计查阅",
        ]),
    ]
    bwid = [0.230, 0.355, 0.355]
    bx = 0.028
    for (title, color, lines), w in zip(branches, bwid):
        _box(ax, bx, dy, w, dh, face=CARD, edge=EDGE, lw=1.1)
        ax.add_patch(Rectangle((bx, dy), 0.006, dh, facecolor=color,
                               edgecolor="none", zorder=3))
        _text(ax, bx + 0.018, dy + dh - 0.022, title, size=8.6, color=color, weight="bold")
        ty = dy + dh - 0.062
        for ln in lines:
            _text(ax, bx + 0.018, ty, ln, size=6.9, color="#2c3e50")
            ty -= 0.036
        bx += w + 0.014

    # ---------- Footnote ----------
    _text(ax, 0.028, 0.040,
          "注：现行取样频次以该系统长期历史数据为支撑（成熟系统）；新系统投用初期通常需按更高频次验证后方可下调。"
          "PW 撬本体运行（加盐再生、滤芯更换、热消毒等）由另外的 SOP 单独管理，不在本流程范围内。",
          size=6.4, color=GREY)
    _text(ax, 0.972, 0.040, "凯莱英 UK Sandwich · B902 Pilot Plant",
          size=6.4, color=GREY, ha="right")

    OUT_PNG.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(OUT_PNG, facecolor=LIGHT)
    plt.close(fig)
    print(f"Wrote {OUT_PNG}")


def build_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(OUT_PNG), 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(OUT_PPTX)
    print(f"Wrote {OUT_PPTX}")


if __name__ == "__main__":
    build_png()
    build_pptx()
