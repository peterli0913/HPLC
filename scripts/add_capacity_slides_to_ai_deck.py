#!/usr/bin/env python3
"""Add Part 1 sub-section slides (HPLC/lyoph + 902 extension) to AI mid-year deck."""

from __future__ import annotations

import io
from copy import deepcopy
from datetime import datetime
from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import font_manager
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SRC = Path("/workspace/AI_推进汇报_2026上半年_v4.pptx")
OUT = Path("/workspace/AI_推进汇报_2026上半年_v4.pptx")
CHART_DIR = Path("/workspace/汇报/_chart_cache")
INSERT_AFTER = 5  # 0-based index — after slide 6 (连续氢化最后一页)

# Theme aligned with existing v4 deck
NAVY = RGBColor(0x00, 0x36, 0x6A)
BODY = RGBColor(0x0F, 0x2B, 0x46)
ORANGE = RGBColor(0xFF, 0x7E, 0x3D)
KPI_BG = RGBColor(0xE7, 0xEE, 0xFB)
TEXT = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY_LINE = RGBColor(0xD4, 0xDD, 0xE6)
EXT_COLOR = "#1a4a6e"
HPLC_COLOR = "#5b6eae"
TEAL = "#009688"
ACCENT = "#c9a227"

FONT = "Microsoft YaHei"
FONT_FALLBACK = "WenQuanYi Micro Hei"

# Data from feasibility briefings (May 2026)
HPLC_BASE = 3_783_000
HPLC_OOM = round(HPLC_BASE * 1.5)  # £5.67M
EXT_OOM = 78_108_089
EXT_BASE = 47_591_969
EXT_RISK = 26_453_763
EXT_GIFA = 3_099


def _chart_font() -> str:
    for name in (FONT, FONT_FALLBACK, "DejaVu Sans"):
        if any(f.name == name for f in font_manager.fontManager.ttflist):
            return name
    return "DejaVu Sans"


def _fill(shape, color: RGBColor, line: RGBColor | None = None):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()


def _tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _run(p, text, *, size=Pt(10), bold=False, color=BODY, font=FONT):
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = size
    r.font.bold = bold
    r.font.color.rgb = color
    return r


def _para(tf, text, *, size=Pt(10), bold=False, color=BODY, align=PP_ALIGN.LEFT, space_after=4):
    p = tf.paragraphs[0] if not tf.text else tf.add_paragraph()
    p.text = text
    p.font.name = FONT
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = Pt(space_after)
    return p


def _bullets(tf, lines, *, size=Pt(8.5), color=BODY):
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = Pt(5)
        p.level = 0


def _slide_footer(slide, page_num: int):
    _para(_tb(slide, Inches(0.30), Inches(5.35), Inches(2.30), Inches(0.20)).text_frame,
          "www.asymchem.com.cn", size=Pt(7), color=TEXT)
    _para(_tb(slide, Inches(2.60), Inches(5.35), Inches(2.50), Inches(0.20)).text_frame,
          "Stock Code: 002821.SZ / 6821.HK", size=Pt(7), color=TEXT)
    _para(_tb(slide, Inches(9.40), Inches(5.35), Inches(0.50), Inches(0.20)).text_frame,
          str(page_num), size=Pt(7), color=TEXT, align=PP_ALIGN.RIGHT)


def _slide_header(slide, title: str, subtitle: str | None = None):
    _para(_tb(slide, Inches(3.50), Inches(0.20), Inches(6.30), Inches(0.50)).text_frame,
          title, size=Pt(20), bold=True, color=NAVY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.30), Inches(0.85), Inches(9.40), Inches(0.03))
    _fill(bar, ORANGE)
    if subtitle:
        _para(_tb(slide, Inches(0.30), Inches(1.00), Inches(9.40), Inches(0.36)).text_frame,
              subtitle, size=Pt(9.5), color=TEXT)


def _kpi_row(slide, kpis: list[tuple[str, str]], top=Inches(1.00)):
    """kpis: [(big_number, label), ...] — up to 4 boxes."""
    n = len(kpis)
    w = Inches(2.30)
    gap = Inches(0.10)
    x0 = Inches(0.30)
    for i, (val, lbl) in enumerate(kpis):
        left = x0 + i * (w + gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(1.15))
        _fill(box, KPI_BG)
        box.line.fill.background()
        _para(_tb(slide, left, top + Inches(0.10), w, Inches(0.55)).text_frame,
              val, size=Pt(22), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        _para(_tb(slide, left, top + Inches(0.70), w, Inches(0.40)).text_frame,
              lbl, size=Pt(9), color=TEXT, align=PP_ALIGN.CENTER)


def _quad_cards(slide, cards: list[tuple[str, list[str]]], top=Inches(2.35)):
    """Four numbered cards in a row."""
    w = Inches(2.30)
    gap = Inches(0.10)
    x0 = Inches(0.30)
    hdr_h = Inches(0.35)
    body_h = Inches(2.15)
    for i, (hdr, bullets) in enumerate(cards):
        left = x0 + i * (w + gap)
        hdr_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, hdr_h)
        _fill(hdr_box, NAVY)
        _para(_tb(slide, left + Inches(0.05), top + Inches(0.02), w - Inches(0.10), hdr_h).text_frame,
              hdr, size=Pt(10.4), bold=True, color=WHITE)
        body_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + hdr_h, w, body_h)
        _fill(body_box, WHITE, GREY_LINE)
        _bullets(_tb(slide, left + Inches(0.08), top + hdr_h + Inches(0.05),
                     w - Inches(0.16), body_h - Inches(0.10)).text_frame, bullets)


def _add_picture(slide, path: Path, left, top, width, height=None):
    if height:
        return slide.shapes.add_picture(str(path), left, top, width, height)
    return slide.shapes.add_picture(str(path), left, top, width=width)


def _fmt_m(n: int) -> str:
    return f"£{n / 1e6:.2f}M"


def _chart_hplc_capex(path: Path):
    cf = _chart_font()
    plt.rcParams["font.family"] = cf
    fig, ax = plt.subplots(figsize=(5.2, 3.6), dpi=150)
    labels = ["直接工程\nA–E", "FEED", "详设", "CDM", "调试", "预备金", "风险预备费"]
    direct = 2_329_000
    other = HPLC_BASE - direct - 631_000
    risk = HPLC_OOM - HPLC_BASE
    vals = [direct, 225_000, 289_000, 235_000, 76_000, 631_000, risk]
    colors = [HPLC_COLOR, "#7b8fc7", "#9aabd4", "#b0bddf", "#c5cfe8", "#d9e0f0", ACCENT]
    bars = ax.barh(labels, [v / 1e6 for v in vals], color=colors, height=0.62)
    ax.set_xlabel("百万英镑 (£M)", fontsize=9)
    ax.set_title("HPLC/冻干改造 · OOM 投资构成", fontsize=11, fontweight="bold", color="#00366a", pad=10)
    ax.axvline(HPLC_OOM / 1e6, color="#b43a2a", linestyle="--", linewidth=1, alpha=0.7)
    ax.text(HPLC_OOM / 1e6 + 0.05, len(labels) - 0.5, f"OOM {_fmt_m(HPLC_OOM)}", fontsize=8, color="#b43a2a")
    for b, v in zip(bars, vals):
        if v > 200_000:
            ax.text(b.get_width() + 0.03, b.get_y() + b.get_height() / 2,
                    f"£{v/1e6:.2f}M", va="center", fontsize=7.5)
    ax.set_xlim(0, max(v / 1e6 for v in vals) * 1.35)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _chart_hplc_gantt(path: Path):
    cf = _chart_font()
    plt.rcParams["font.family"] = cf
    fig, ax = plt.subplots(figsize=(5.2, 3.6), dpi=150)
    rows = [
        ("FS/基准", "2026-05-01", "2026-05-31", TEAL),
        ("冻干规格/资金", "2026-05-26", "2026-07-20", ACCENT),
        ("冻干制造", "2026-07-28", "2027-04-05", ACCENT),
        ("冻干 FAT→PQ", "2027-04-06", "2027-11-08", ACCENT),
        ("FEED（示意）", "2026-06-15", "2026-09-20", HPLC_COLOR),
        ("HPLC 制造", "2026-10-20", "2027-02-22", "#1a4a6e"),
        ("HPLC 安装验证", "2027-03-30", "2027-09-30", "#1a4a6e"),
        ("厂房改造", "2026-12-24", "2027-04-15", "#2e6da4"),
    ]
    t0 = datetime(2026, 4, 1)
    t1 = datetime(2028, 1, 1)
    span = (t1 - t0).days

    def _x(d):
        return (datetime.strptime(d, "%Y-%m-%d") - t0).days / span

    for i, (lbl, s, e, c) in enumerate(rows):
        ax.barh(i, _x(e) - _x(s), left=_x(s), height=0.55, color=c, alpha=0.9)
        ax.text(-0.02, i, lbl, ha="right", va="center", fontsize=7.5)

    today = _x("2026-06-30")
    ax.axvline(today, color="#b43a2a", linewidth=1.5, linestyle="--", alpha=0.8)
    ax.text(today, len(rows) - 0.3, " 约今", fontsize=7, color="#b43a2a")

    ax.set_yticks([])
    ax.set_xticks([_x(f"{y}-01-01") for y in range(2026, 2029)])
    ax.set_xticklabels(["2026", "2027", "2028"], fontsize=8)
    ax.set_xlim(-0.22, 1.02)
    ax.set_title("改造 · 关键路径示意（冻干长周期驱动）", fontsize=11, fontweight="bold", color="#00366a", pad=8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _chart_ext_capex(path: Path):
    cf = _chart_font()
    plt.rcParams["font.family"] = cf
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(5.4, 3.6), dpi=150)

    # Left: OOM composition
    labels = ["基础建造", "其他项目费", "风险与预备费"]
    vals = [EXT_BASE, EXT_OOM - EXT_BASE - EXT_RISK, EXT_RISK]
    colors = [EXT_COLOR, TEAL, ACCENT]
    wedges, _, autotexts = ax1.pie(
        [v / 1e6 for v in vals], labels=labels, autopct="%1.0f%%",
        colors=colors, startangle=90, textprops={"fontsize": 7.5},
    )
    for at in autotexts:
        at.set_fontsize(7)
    ax1.set_title(f"OOM {_fmt_m(EXT_OOM)}", fontsize=10, fontweight="bold", color="#00366a")

    # Right: base construction split
    sub_labels = ["土建外围", "建筑机电", "工艺设备", "临建+OH&P"]
  # approximate from briefing
    sub_vals = [12.5, 11.8, 23.25, 0.0]  # M£ — process dominant
    # Use stacked from known: works subtotal ~40M + prelims + ohp
    works = [8.2, 7.9, 23.25]
    wlabels = ["土建", "机电", "工艺设备"]
    ax2.bar(wlabels, works, color=[EXT_COLOR, TEAL, HPLC_COLOR], width=0.55)
    ax2.set_ylabel("£M", fontsize=8)
    ax2.set_title("基础建造费分项（示意）", fontsize=9, fontweight="bold", color="#00366a")
    for i, v in enumerate(works):
        ax2.text(i, v + 0.4, f"£{v:.1f}M", ha="center", fontsize=7.5)
    ax2.set_ylim(0, 28)
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)

    fig.suptitle("902 东侧扩建 · 投资结构", fontsize=11, fontweight="bold", color="#00366a", y=1.02)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _chart_ext_gantt(path: Path):
    cf = _chart_font()
    plt.rcParams["font.family"] = cf
    fig, ax = plt.subplots(figsize=(5.4, 3.6), dpi=150)
    rows = [
        ("RIBA 1 可行性", "2026-05-01", "2026-05-31", TEAL),
        ("RIBA 2 概念", "2026-07-15", "2027-02-23", EXT_COLOR),
        ("RIBA 3 方案", "2027-02-17", "2027-09-28", EXT_COLOR),
        ("规划审批", "2027-10-26", "2028-03-21", ACCENT),
        ("RIBA 4 详设", "2027-10-27", "2028-06-06", "#5b6eae"),
        ("长周期设备", "2028-03-08", "2029-07-24", ACCENT),
        ("RIBA 5 施工", "2028-08-02", "2029-11-13", "#2e6da4"),
        ("调试/竣工", "2029-11-14", "2030-05-07", "#2e6da4"),
    ]
    t0 = datetime(2026, 1, 1)
    t1 = datetime(2031, 1, 1)
    span = (t1 - t0).days

    def _x(d):
        return (datetime.strptime(d, "%Y-%m-%d") - t0).days / span

    for i, (lbl, s, e, c) in enumerate(rows):
        ax.barh(i, _x(e) - _x(s), left=_x(s), height=0.55, color=c, alpha=0.9)
        ax.text(-0.02, i, lbl, ha="right", va="center", fontsize=7.2)

    today = _x("2026-06-30")
    ax.axvline(today, color="#b43a2a", linewidth=1.5, linestyle="--", alpha=0.8)
    ax.text(today, len(rows) - 0.2, " 约今", fontsize=7, color="#b43a2a")

    ax.set_yticks([])
    ax.set_xticks([_x(f"{y}-01-01") for y in range(2026, 2031)])
    ax.set_xticklabels([str(y) for y in range(2026, 2031)], fontsize=7)
    ax.set_xlim(-0.25, 1.02)
    ax.set_title("扩建 · 总控节奏（1,005 天）", fontsize=11, fontweight="bold", color="#00366a", pad=8)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_visible(False)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _chart_portfolio_overview(path: Path):
    cf = _chart_font()
    plt.rcParams["font.family"] = cf
    fig, ax = plt.subplots(figsize=(4.8, 3.2), dpi=150)
    projects = ["HPLC/冻干\n改造", "902 东侧\n扩建"]
    oom = [HPLC_OOM / 1e6, EXT_OOM / 1e6]
    bars = ax.bar(projects, oom, color=[HPLC_COLOR, EXT_COLOR], width=0.5)
    ax.set_ylabel("OOM（£M）", fontsize=9)
    ax.set_title("双轨投资量级对比（可行性 OOM）", fontsize=10, fontweight="bold", color="#00366a")
    for b, v in zip(bars, oom):
        ax.text(b.get_x() + b.get_width() / 2, v + 1.5, f"£{v:.1f}M", ha="center", fontsize=9, fontweight="bold")
    ax.set_ylim(0, max(oom) * 1.15)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    # milestone annotations
    ax.annotate("2027 近期能力", xy=(0, oom[0]), xytext=(0.3, oom[0] + 8),
                fontsize=7, color=HPLC_COLOR, arrowprops=dict(arrowstyle="->", color=HPLC_COLOR, lw=0.8))
    ax.annotate("2030 中长期", xy=(1, oom[1]), xytext=(0.6, oom[1] - 12),
                fontsize=7, color=EXT_COLOR, arrowprops=dict(arrowstyle="->", color=EXT_COLOR, lw=0.8))
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def _draw_site_schematic(slide, left, top, width, height):
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(outer, WHITE, GREY_LINE)
    _para(_tb(slide, left + Inches(0.10), top + Inches(0.08), width - Inches(0.20), Inches(0.28)).text_frame,
          "Sandwich 园区 · B902 产能建设示意", size=Pt(9.5), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    exist = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top + Inches(0.45),
                                   Inches(1.55), Inches(1.05))
    _fill(exist, KPI_BG, RGBColor(0x9A, 0xA8, 0xB6))
    _para(_tb(slide, left + Inches(0.20), top + Inches(0.62), Inches(1.25), Inches(0.75)).text_frame,
          "既有 B902\n色谱/冻干/密闭升级", size=Pt(8), color=BODY, align=PP_ALIGN.CENTER)

    ext = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(1.80), top + Inches(0.45),
                                 Inches(1.15), Inches(1.05))
    _fill(ext, RGBColor(0xD9, 0xEA, 0xE8), RGBColor(0x00, 0x96, 0x88))
    _para(_tb(slide, left + Inches(1.65), top + Inches(0.68), Inches(0.95), Inches(0.55)).text_frame,
          "东侧扩建\n反应/加氢/过滤", size=Pt(8), color=RGBColor(0x00, 0x96, 0x88), align=PP_ALIGN.CENTER)

    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(0.55), top + Inches(1.58),
                                 Inches(2.10), Inches(0.22))
    _fill(arr, RGBColor(0xE8, 0xEE, 0xF4), GREY_LINE)
    _para(_tb(slide, left + Inches(0.15), top + Inches(1.88), width - Inches(0.30), Inches(0.35)).text_frame,
          "近期改造兑现 2027 能力  →  扩建 2030 规模化产能", size=Pt(7.5), color=TEXT, align=PP_ALIGN.CENTER)


def _insert_slide(prs: Presentation, index: int):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(index, slides[-1])
    return prs.slides[index]


def _build_slide_overview(slide, charts: dict, page: int):
    _slide_header(
        slide,
        "B902 产能建设 · 双轨并行总览",
        "既有厂房改造（近期能力）+ 东侧扩建（中长期产能）· 我作为国内 PM 支持统筹信息闭环与决策材料",
    )
    _kpi_row(slide, [
        ("双轨 FS", "扩建 + 改造 可行性已完成"),
        ("投资测算", "OOM 框架已明确"),
        ("多版材料", "组合汇报支撑决策"),
        ("接口统筹", "改造/扩建/高活/C1 联动"),
    ], top=Inches(1.45))

    _draw_site_schematic(slide, Inches(0.30), Inches(2.70), Inches(3.20), Inches(2.30))
    _add_picture(slide, charts["portfolio"], Inches(3.70), Inches(2.70), Inches(3.10), Inches(2.30))

    # Comparison table
    tbl_top = Inches(2.70)
    tbl_left = Inches(6.95)
    headers = ["维度", "HPLC/冻干改造", "902 东侧扩建"]
    rows = [
        ["定位", "近期投产能力", "中长期规模化产能"],
        ["投资 OOM", _fmt_m(HPLC_OOM), _fmt_m(EXT_OOM)],
        ["关键节点", "2027 Q3–Q4", "2030-05 竣工"],
        ["上半年", "FS + 投资框架", "RIBA 1 完成（5月）"],
        ["下半年", "FEED / 长周期采购", "RIBA 2 概念设计"],
    ]
    ts = slide.shapes.add_table(len(rows) + 1, 3, tbl_left, tbl_top, Inches(2.75), Inches(2.30))
    table = ts.table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT
                p.font.size = Pt(7.5)
                p.font.color.rgb = BODY if j else TEXT
                p.font.bold = j == 0
    _slide_footer(slide, page)


def _build_slide_hplc_progress(slide, page: int):
    _slide_header(
        slide,
        "制备色谱 & 冻干机改造 · 上半年推进主线",
        "技术可行性研究已完成；投资框架明确；冻干机长周期驱动分期实施",
    )
    _kpi_row(slide, [
        ("FS 完成", "可行性研究 P01"),
        (_fmt_m(HPLC_OOM), "OOM（可行性量级）"),
        ("2027-09", "HPLC 目标投运"),
        ("2027-12", "冻干目标投运"),
    ], top=Inches(1.45))
    _quad_cards(slide, [
        ("① 可行性闭环", [
            "• 全程跟进英国侧 FS 讲解及多轮专题会",
            "• 方案布局与投资口径中英方一致理解",
            "• 可行性结论及投资框架已明确",
        ]),
        ("② 投资与汇报", [
            "• 汇总投资与进度，编制多版组合汇报材料",
            "• 支撑管理层决策与下阶段沟通",
            "• OOM 基础 £3.78M + 风险预备费块",
        ]),
        ("③ 接口统筹", [
            "• 统筹色谱冻干与 C1 五级密闭升级",
            "• 衔接高活实验室等配套事项",
            "• 避免各子项割裂推进",
        ]),
        ("④ 下半年重点", [
            "• 推动 FEED 启动与详设准备",
            "• 长周期设备（冻干）采购决策",
            "• 中外需求最终确认与常规沟通机制",
        ]),
    ], top=Inches(2.70))
    _slide_footer(slide, page)


def _build_slide_hplc_charts(slide, charts: dict, page: int):
    _slide_header(
        slide,
        "制备色谱 & 冻干机改造 · 投资结构与周期",
        "范围：DAC300/CP300 制备色谱 + 冻干机（隔离器、除湿、PSG）及厂房改造 · RB Plant 9802",
    )
    _add_picture(slide, charts["hplc_capex"], Inches(0.30), Inches(1.45), Inches(4.55), Inches(3.55))
    _add_picture(slide, charts["hplc_gantt"], Inches(5.00), Inches(1.45), Inches(4.70), Inches(3.55))

    # Bottom milestone strip
    milestones = [
        ("Q3 2026", "FEED 启动\n冻干规格/资金"),
        ("Q4 2026", "详设准备\nHPLC 规格/资金"),
        ("2027 H1", "设备 FAT\n厂房改造"),
        ("2027 H2", "HPLC 投运\n冻干验证"),
    ]
    for i, (period, desc) in enumerate(milestones):
        left = Inches(0.30) + i * Inches(2.40)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.05), Inches(2.25), Inches(0.55))
        _fill(box, KPI_BG)
        box.line.color.rgb = RGBColor(0x5B, 0x6E, 0xAE)
        tf = _tb(slide, left + Inches(0.10), Inches(5.08), Inches(2.05), Inches(0.50)).text_frame
        p = tf.paragraphs[0]
        _run(p, period + "  ", size=Pt(8), bold=True, color=NAVY)
        _run(p, desc, size=Pt(7.5), color=BODY)
    _slide_footer(slide, page)


def _build_slide_ext_progress(slide, page: int):
    _slide_header(
        slide,
        "902 东侧扩建 · 上半年推进主线",
        "Scitech 可行性研究（RIBA 1）2026年5月完成；推荐 Option 1 方案；下半年进入概念设计",
    )
    _kpi_row(slide, [
        ("RIBA 1", "2026-05 FS 完成"),
        (_fmt_m(EXT_OOM), "项目总投资 OOM"),
        (f"{EXT_GIFA:,} m²", "总建筑面积 GIFA"),
        ("2030-05", "总控竣工目标"),
    ], top=Inches(1.45))
    _quad_cards(slide, [
        ("① FS 消化闭环", [
            "• 全程参与 FS 讲解会，整理中文详细纪要",
            "• 推动国内团队理解推荐方案及前提假设",
            "• BREEAM 预评估 Very Good 基线",
        ]),
        ("② 投资与决策", [
            "• 完成投资量级、风险与总控计划解读",
            "• 编制管理层汇报材料支撑决策",
            "• 10 釜 + 2500L 加氢釜 + 3 套过滤干燥机",
        ]),
        ("③ 关键事项跟进", [
            "• 加氢厂房处置、首层净高等中外确认项",
            "• 长周期设备采购分工待明确",
            "• 促进信息对称、问题及时闭环",
        ]),
        ("④ 下半年重点", [
            "• RIBA 2 概念设计启动（7月计划）",
            "• 关键方案国内确认",
            "• 设计工作坊与 Stage 2 报价对接",
        ]),
    ], top=Inches(2.70))
    _slide_footer(slide, page)


def _build_slide_ext_charts(slide, charts: dict, page: int):
    _slide_header(
        slide,
        "902 东侧扩建 · 投资结构与总控节奏",
        "东侧约 600 m² 占地 · 四层+设备夹层 · 与 902 低层楼面贯通 · 总控 1,005 天",
    )
    _add_picture(slide, charts["ext_capex"], Inches(0.30), Inches(1.45), Inches(4.55), Inches(3.55))
    _add_picture(slide, charts["ext_gantt"], Inches(5.00), Inches(1.45), Inches(4.70), Inches(3.55))

    decisions = [
        ("H2 2026", "RIBA 2 概念设计\nStage 2 报价确认"),
        ("2027", "方案设计 +\n规划审批"),
        ("2028", "详设 +\n长周期设备采购"),
        ("2028–30", "施工 · 调试 · 竣工"),
    ]
    for i, (period, desc) in enumerate(decisions):
        left = Inches(0.30) + i * Inches(2.40)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.05), Inches(2.25), Inches(0.55))
        _fill(box, KPI_BG)
        box.line.color.rgb = RGBColor(0x1A, 0x4A, 0x6E)
        tf = _tb(slide, left + Inches(0.10), Inches(5.08), Inches(2.05), Inches(0.50)).text_frame
        p = tf.paragraphs[0]
        _run(p, period + "  ", size=Pt(8), bold=True, color=NAVY)
        _run(p, desc, size=Pt(7.5), color=BODY)
    _slide_footer(slide, page)


def _renumber_footers(prs: Presentation):
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            t = shape.text_frame.text.strip()
            if t.isdigit() and shape.left > Emu(int(8.5 * 914400)):
                shape.text_frame.paragraphs[0].text = str(i + 1)


def main():
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    charts = {
        "portfolio": CHART_DIR / "portfolio_overview.png",
        "hplc_capex": CHART_DIR / "hplc_capex.png",
        "hplc_gantt": CHART_DIR / "hplc_gantt.png",
        "ext_capex": CHART_DIR / "ext_capex.png",
        "ext_gantt": CHART_DIR / "ext_gantt.png",
    }
    _chart_portfolio_overview(charts["portfolio"])
    _chart_hplc_capex(charts["hplc_capex"])
    _chart_hplc_gantt(charts["hplc_gantt"])
    _chart_ext_capex(charts["ext_capex"])
    _chart_ext_gantt(charts["ext_gantt"])

    prs = Presentation(str(SRC))
    builders = [
        lambda s, p: _build_slide_overview(s, charts, p),
        lambda s, p: _build_slide_hplc_progress(s, p),
        lambda s, p: _build_slide_hplc_charts(s, charts, p),
        lambda s, p: _build_slide_ext_progress(s, p),
        lambda s, p: _build_slide_ext_charts(s, charts, p),
    ]
    insert_at = INSERT_AFTER + 1
    for i, builder in enumerate(builders):
        slide = _insert_slide(prs, insert_at + i)
        builder(slide, insert_at + i + 1)

    _renumber_footers(prs)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
