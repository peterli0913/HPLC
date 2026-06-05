#!/usr/bin/env python3
"""Generate COO/CFO briefing PPT from verified Scitech feasibility data only."""

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --- Verified data (300291 FS pack, Issue A1, 22-May-2026) ---
DATA_SOURCE = (
    "Scitech-EKIUM: 300291-RE-0001 Feasibility Report, "
    "300291-CM-0001 Cost Plan, 300291-PM-PR-0002 Programme (May 2026)"
)

TOTAL_OOM_GBP = 78_108_089
WORKS_SUBTOTAL = 41_968_226
BASE_CONSTRUCTION = 47_591_969
OTHER_PROJECT = 4_062_357
RISK_CONTINGENCY = 25_827_163
RISK_REGISTER = 626_600

COST_BREAKDOWN = [
    ("Process equipment (works line 5.6)", 23_249_083),
    ("Building services (HVAC/MEP)", 9_532_716),
    ("External works", 1_859_520),
    ("Superstructure", 3_942_636),
    ("Substructure", 1_499_680),
    ("Internal finishes", 1_270_672),
    ("Facilitating / enabling", 254_960),
    ("Works to existing building", 204_000),
    ("Fittings & equipment", 154_960),
]

CONTINGENCY_STACK = [
    ("Design development (15%)", 7_748_149),
    ("Construction & equipment risk (25%)", 12_913_581),
    ("Client contingency (10%)", 5_165_433),
    ("Risk register (factored)", 626_600),
]

OTHER_COSTS = [
    ("Design team fees (8% of construction)", 3_807_357),
    ("BREEAM", 150_000),
    ("Surveys & investigations", 80_000),
    ("Planning & statutory", 25_000),
]

GIFA_M2 = 3_099
COST_PER_M2 = round(TOTAL_OOM_GBP / GIFA_M2)

MILESTONES = [
    ("Feasibility study issued (RIBA 1)", "May 2026", "Complete"),
    ("Concept design (RIBA 2)", "Jul 2026 – Feb 2027", "Planned"),
    ("Scheme design (RIBA 3)", "Feb 2027 – Sep 2027", "Planned"),
    ("Planning consent (target)", "Mar 2028", "Planned"),
    ("Detailed design (RIBA 4)", "Oct 2027 – Jun 2028", "Planned"),
    ("Construction (RIBA 5)", "Aug 2028 – Nov 2029", "Planned"),
    ("Project complete (programme)", "May 2030", "Planned"),
]

NAVY = RGBColor(0x0F, 0x2B, 0x46)
TEAL = RGBColor(0x00, 0x96, 0x88)
GREY = RGBColor(0x5A, 0x6A, 0x7A)


def fmt_gbp(n: int) -> str:
    return f"£{n:,}"


def set_title_style(shape):
    if not shape.has_text_frame:
        return
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            r.font.name = "Arial"
            r.font.bold = True
            r.font.color.rgb = NAVY


def add_footer(slide, note: str = ""):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(9.2), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Source: {DATA_SOURCE}  |  {note}" if note else f"Source: {DATA_SOURCE}"
    p.font.size = Pt(8)
    p.font.color.rgb = GREY


def slide_title(prs, title, subtitle=""):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # header bar
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY

    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(2.5))
        stf = sbox.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = GREY
    add_footer(slide)
    return slide


def slide_bullets(prs, title, bullets, subtitle=""):
    slide = slide_title(prs, title, subtitle)
    box = slide.shapes.add_textbox(Inches(0.55), Inches(2.0 if subtitle else 1.5), Inches(9), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(8)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    return slide


def slide_cost_waterfall(prs):
    slide = slide_title(prs, "Investment Structure (OOM)", "Feasibility-stage order-of-magnitude — not final Capex")

    # Stacked bar data: Works | Prelims+OH&P uplift to base | Other | Risk&Cont
    prelim_ohp = BASE_CONSTRUCTION - WORKS_SUBTOTAL
    chart_data = CategoryChartData()
    chart_data.categories = ["Published OOM total"]
    chart_data.add_series("Direct works", (WORKS_SUBTOTAL,))
    chart_data.add_series("Preliminaries & OH&P", (prelim_ohp,))
    chart_data.add_series("Other project costs", (OTHER_PROJECT,))
    chart_data.add_series("Risk & contingency", (RISK_CONTINGENCY + RISK_REGISTER,))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(0.5), Inches(1.8), Inches(5.5), Inches(4.2), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.include_in_layout = False

    # table right
    rows = [
        ("Total OOM", fmt_gbp(TOTAL_OOM_GBP)),
        ("Base construction cost", fmt_gbp(BASE_CONSTRUCTION)),
        ("  → Direct works subtotal", fmt_gbp(WORKS_SUBTOTAL)),
        ("Other project costs", fmt_gbp(OTHER_PROJECT)),
        ("Risk & contingency (incl. register)", fmt_gbp(RISK_CONTINGENCY + RISK_REGISTER)),
        ("Unit cost (GIFA)", f"{fmt_gbp(COST_PER_M2)}/m²"),
        ("GIFA (feasibility)", f"{GIFA_M2:,} m²"),
    ]
    tbl = slide.shapes.add_table(len(rows), 2, Inches(6.2), Inches(2.0), Inches(3.3), Inches(3.5)).table
    for i, (a, b) in enumerate(rows):
        tbl.cell(i, 0).text = a
        tbl.cell(i, 1).text = b
        for c in range(2):
            for p in tbl.cell(i, c).text_frame.paragraphs:
                p.font.size = Pt(10)

    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.7))
    note.text_frame.paragraphs[0].text = (
        "Excludes: inflation/FX/tax; client PM & operational costs. "
        "Cost plan states client-purchased equipment excluded — reconcile £23.25M process line at Concept."
    )
    note.text_frame.paragraphs[0].font.size = Pt(9)
    note.text_frame.paragraphs[0].font.color.rgb = GREY
    return slide


def slide_cost_pie(prs):
    slide = slide_title(prs, "Direct Works Breakdown", f"Works subtotal {fmt_gbp(WORKS_SUBTOTAL)}")

    chart_data = CategoryChartData()
    chart_data.categories = [x[0][:28] for x in COST_BREAKDOWN]
    chart_data.add_series("GBP", tuple(x[1] for x in COST_BREAKDOWN))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(0.4), Inches(1.7), Inches(5.8), Inches(4.5), chart_data
    ).chart
    chart.has_legend = True

    lines = [f"{n}: {fmt_gbp(v)} ({v/WORKS_SUBTOTAL*100:.1f}%)" for n, v in COST_BREAKDOWN[:6]]
    box = slide.shapes.add_textbox(Inches(6.3), Inches(2.0), Inches(3.4), Inches(4))
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
    return slide


def slide_timeline(prs):
    slide = slide_title(prs, "Indicative Programme (Scitech Appendix 8)")

    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5))
    tf = box.text_frame
    for i, (name, period, status) in enumerate(MILESTONES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{name}:  {period}  [{status}]"
        p.font.size = Pt(13)
        p.space_after = Pt(6)

    warn = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(0.9))
    warn.text_frame.paragraphs[0].text = (
        "Risk register #13: desired Q3 2028 completion noted as tight vs master programme "
        "(construction through Nov 2029; project complete May 2030). Milestone definitions require alignment."
    )
    warn.text_frame.paragraphs[0].font.size = Pt(10)
    warn.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xB4, 0x3A, 0x2A)
    return slide


def build_ppt(path: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(
        prs,
        "Asymchem UK — B902 PDF Extension",
        "Investment & Programme Briefing for Group COO/CFO\n"
        "Feasibility Study (RIBA Stage 1)  |  Discovery Park, Sandwich  |  May 2026",
    )

    slide_bullets(
        prs,
        "Executive Summary",
        [
            f"Conclusion: Extension is feasible subject to planning approval (Scitech executive summary).",
            f"Investment (OOM, May 2026 base): {fmt_gbp(TOTAL_OOM_GBP)} total project cost; "
            f"base construction {fmt_gbp(BASE_CONSTRUCTION)}.",
            f"Scale: {GIFA_M2:,} m² GIFA (~600 m² footprint, 4 floors + plant mezzanine).",
            "Schedule: Feasibility complete; Concept design planned from Jul 2026; "
            "construction Aug 2028–Nov 2029; project complete May 2030 per programme.",
            "Decision: Approve progression to RIBA Stage 2 and confirm baseline Option 1 "
            "(hydrogenation building removal).",
        ],
    )

    slide_bullets(
        prs,
        "Project Overview",
        [
            "Asset: Existing B902 PDF pilot plant — east-side extension for reactor scale-up & new suites.",
            "Scope: 10 reactors (R25–R34), 2,500 L hydrogenation reactor (R35), 3 filter dryers, "
            "handling/milling/continuous areas, HVAC & utilities.",
            "Designer: Scitech-EKIUM under contract to Asymchem (Project 300291).",
            "Baseline design: Option 1 — remove/relocate existing hydrogenation building; "
            "highest score in option appraisal matrix.",
        ],
    )

    slide_bullets(
        prs,
        "Current Progress",
        [
            "RIBA Stage 1 Feasibility Report & appendices issued (Rev A1, 22 May 2026).",
            "Feasibility walkthrough with UK team: 27 May 2026.",
            "Pre-planning pack in preparation for local authority engagement.",
            "Outstanding: B902 operating data gaps; final confirmation of procurement split; "
            "Stage 2 contract / fee proposal (not in published FS pack).",
        ],
    )

    slide_cost_waterfall(prs)
    slide_cost_pie(prs)

    slide_bullets(
        prs,
        "Contingency & Risk Allowance (Feasibility)",
        [f"{n}: {fmt_gbp(v)}" for n, v in CONTINGENCY_STACK]
        + [
            "Combined risk & contingency provisions are material (~33% on top of base construction "
            "plus other costs) — appropriate for Stage 1 OOM, subject to drawdown governance.",
            f"Pre-mitigation risk exposure (report §14.3): ~£2.64m cost / 182 weeks programme (factored lower).",
        ],
    )

    slide_bullets(
        prs,
        "Investment Analysis — CFO Considerations",
        [
            "Nature of estimate: Order-of-magnitude only; not for final Capex approval (cost plan disclaimer).",
            "No escalation, inflation, FX, or tax in base date 22-May-2026.",
            "Client-purchased process equipment stated excluded; process line £23.25M in works "
            "schedule requires reconciliation at Concept Design to avoid double-count or gap.",
            "Excluded: client project management, operations, consumables, survey remediation beyond allowances.",
            f"Design fees included: {fmt_gbp(3_807_357)} (8% of construction) — whole-project allowance, not FEED-only.",
            "Accuracy improves at Concept (RIBA 2) with surveys, procurement strategy, and equipment quotes.",
        ],
    )

    slide_timeline(prs)

    slide_bullets(
        prs,
        "Critical Decisions & Next 90 Days",
        [
            "Approve Stage 2 Concept Design appointment (planned 1 Jul 2026).",
            "Confirm Option 1 — hydrogenation building demolition/relocation.",
            "Resolve ground-floor Hastelloy reactor headroom strategy (lower slab vs relocate).",
            "Align internal completion target (Q3 2028 aspiration vs May 2030 programme).",
            "Authorize pre-planning submission and long-lead equipment planning in Concept phase.",
        ],
        subtitle="Gate: Concept Stage Gate instruction to proceed — 16 Feb 2027 (programme)",
    )

    slide_bullets(
        prs,
        "Recommendation to Group Leadership",
        [
            "Support conditional advancement to RIBA Stage 2 with planning consent as external gate.",
            "Budget planning: use £78.1M OOM as envelope for dialogue; hold separate provision for "
            "client-direct equipment and inflation until Concept cost plan issued.",
            "Mandate procurement strategy workshop (risk #18) before major commitments.",
            "Request Scitech stage-gated fee quote and cash-flow phasing for FY planning.",
        ],
    )

    slide_bullets(
        prs,
        "Data Integrity & Limitations",
        [
            "All figures from published Scitech feasibility deliverables in repository (May 2026).",
            "No RMB conversion or NPV presented — requires group FX and hurdle-rate policy.",
            "Programme dates are indicative; Stage 2+ contract not included in this dataset.",
            "For audit trail: 300291-RE-0001, 300291-CM-0001, 300291-PM-PR-0002, 300291-PM-RA-0001.",
        ],
    )

    prs.save(path)
    print(f"Saved {path}")


if __name__ == "__main__":
    out = Path(__file__).parent / "PDF-Extension_COO-CFO_Briefing_2026-05-28.pptx"
    build_ppt(out)
